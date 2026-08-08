from pathlib import Path

from pptx import Presentation

from app.parsers.base_parser import BaseParser


class PPTXParser(BaseParser):
    """
    Parser for PowerPoint PPTX documents.
    """

    def parse(self, file_path: Path) -> str:
        presentation = Presentation(str(file_path))

        slides_text = []

        for slide in presentation.slides:
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slide_text.append(text)

            if slide_text:
                slides_text.append(
                    "\n".join(slide_text)
                )

        return "\n\n".join(slides_text)